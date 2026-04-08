import sys
from pptx import Presentation

def inspect_pptx(file_path):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error opening PowerPoint: {e}")
        return

    print(f"Analyzing: {file_path}")
    print(f"Total Slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
            except Exception:
                left = top = width = height = "N/A"

            # Determine shape type and content
            shape_info = ""
            if shape.has_text_frame:
                text = shape.text.replace("\n", " ").replace("\r", "").strip()
                text = (text[:50] + "...") if len(text) > 50 else text
                shape_info = f"[TEXT]: {text}"
            elif shape.has_chart:
                shape_info = f"[CHART]: {shape.chart.chart_type}"
            elif shape.has_table:
                shape_info = "[TABLE]"
            elif getattr(shape, "shape_type", None) == 13: # msoPicture
                shape_info = "[PICTURE]"
            else:
                shape_type = getattr(shape, "shape_type", "Unknown")
                shape_info = f"[OTHER]: Type {shape_type}"

            print(f"  Shape: {shape.name}")
            if left != "N/A":
                print(f"    Pos: Left={left}, Top={top}, W={width}, H={height}")
            print(f"    Data: {shape_info}")
        print()

if __name__ == "__main__":
    pptx_path = "Merkle Thailand -Ajipanda's Kitchen report- 260331  copy.pptx"
    inspect_pptx(pptx_path)
