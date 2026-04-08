import sys
import os
import json
import pandas as pd
from pptx import Presentation

def extract_pptx_text(pptx_path):
    print(f"Extracting text from: {pptx_path}")
    prs = Presentation(pptx_path)
    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides_data.append(f"Slide {i+1}: " + " | ".join(slide_text).replace('\n', ' '))
    
    print("--- PPTX KEY FINDINGS ---")
    for text in slides_data:
        print(text)
    print("-" * 25)

def extract_xlsx_data(xlsx_path):
    print(f"Extracting data from: {xlsx_path}")
    xls = pd.ExcelFile(xlsx_path)
    print("--- XLSX DATA SUMMARY ---")
    for sheet_name in xls.sheet_names:
        print(f"Sheet: {sheet_name}")
        try:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            # Dump the first 20 rows of each sheet in CSV format for compactness
            csv_data = df.head(20).to_csv(index=False)
            print(csv_data)
        except Exception as e:
            print(f"Error reading sheet {sheet_name}: {e}")
        print("-" * 25)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_data.py <pptx_path> <xlsx_path>")
        sys.exit(1)
        
    pptx_path = sys.argv[1]
    xlsx_path = sys.argv[2]
    
    if not os.path.exists(pptx_path):
        print(f"Error: {pptx_path} not found.")
        sys.exit(1)
    if not os.path.exists(xlsx_path):
        print(f"Error: {xlsx_path} not found.")
        sys.exit(1)
        
    extract_pptx_text(pptx_path)
    extract_xlsx_data(xlsx_path)
