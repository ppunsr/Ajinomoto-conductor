import sys
import json
import argparse
from pptx import Presentation

def inspect_pptx_to_json(file_path, output_json_path):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error opening PowerPoint: {e}")
        sys.exit(1)

    print(f"Analyzing: {file_path}")
    
    data = {
        "file": file_path,
        "total_slides": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i + 1,
            "shapes": []
        }
        for shape in slide.shapes:
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
            except Exception:
                left = top = width = height = None

            shape_info = {}
            if shape.has_text_frame:
                text = shape.text.replace("\n", " ").replace("\r", "").strip()
                shape_info["type"] = "TEXT"
                shape_info["content"] = text
            elif shape.has_chart:
                shape_info["type"] = "CHART"
                shape_info["chart_type"] = str(shape.chart.chart_type)
            elif shape.has_table:
                shape_info["type"] = "TABLE"
            elif getattr(shape, "shape_type", None) == 13: # msoPicture
                shape_info["type"] = "PICTURE"
            else:
                shape_type = getattr(shape, "shape_type", "Unknown")
                shape_info["type"] = f"OTHER_{shape_type}"

            shape_data = {
                "name": shape.name,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height
                },
                "data": shape_info
            }
            slide_data["shapes"].append(shape_data)
        
        data["slides"].append(slide_data)

    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=4)
        
    print(f"Saved layout data to: {output_json_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract layout mapping from a PowerPoint file.")
    parser.add_argument("input", help="Path to the input PPTX file")
    parser.add_argument("output", help="Path to the output JSON file")
    args = parser.parse_args()
    
    inspect_pptx_to_json(args.input, args.output)
